#!/usr/bin/env python3
from __future__ import annotations

import argparse
import math
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400


FORBIDDEN_FONTS = {
    "Georgia",
    "Times",
    "Times New Roman",
    "Cambria",
    "Garamond",
    "Baskerville",
}

FORBIDDEN_COLOR_HEXES = {
    "F6F2E8",
    "F7F4EE",
    "FAF8F3",
    "F8F5EE",
    "EEE6D7",
    "EFE7D7",
}

PLACEHOLDER_PATTERNS = (
    "TODO",
    "TBD",
    "LOREM",
    "PLACEHOLDER",
)

INTERNAL_WORKFLOW_TEXT_PATTERNS = (
    re.compile(r"\bsource[-\s]?log\b", re.IGNORECASE),
    re.compile(r"\bproject\s+source\s+log\b", re.IGNORECASE),
    re.compile(r"\bfinal[-\s]?audit\b", re.IGNORECASE),
    re.compile(r"\bworking[-\s]?audit\b", re.IGNORECASE),
    re.compile(r"\brequirement[-\s]?ledger\b", re.IGNORECASE),
    re.compile(r"\brunpro_workspace\b", re.IGNORECASE),
    re.compile(r"\bpptpro[_\s-]?audit\b", re.IGNORECASE),
    re.compile(r"\bfull\s+source\s+provenance\b", re.IGNORECASE),
    re.compile(r"\bimage\s+credits?\s+(?:are|is)\s+recorded\b", re.IGNORECASE),
    re.compile(r"\b(?:source|image|generated[-\s]?image)\s+provenance\b.{0,80}\brecorded\b", re.IGNORECASE),
    re.compile(r"\bprovenance\b.{0,80}\b(?:project\s+source\s+log|source[-\s]?log)\b", re.IGNORECASE),
    re.compile(r"\brecorded\b.{0,80}\b(?:source[-\s]?log|final[-\s]?audit|working[-\s]?audit)\b", re.IGNORECASE),
    re.compile(r"来源日志|项目来源日志|最终审计|工作审计|需求台账"),
    re.compile(r"记录在.{0,40}(?:来源日志|最终审计|工作审计|项目日志)"),
)

NON_EMPTY_VISUAL_TYPES = {
    value
    for value in (
        getattr(MSO_SHAPE_TYPE, "PICTURE", None),
        getattr(MSO_SHAPE_TYPE, "CHART", None),
        getattr(MSO_SHAPE_TYPE, "TABLE", None),
        getattr(MSO_SHAPE_TYPE, "GROUP", None),
        getattr(MSO_SHAPE_TYPE, "MEDIA", None),
    )
    if value is not None
}

AUTO_SHAPE_TYPE = getattr(MSO_SHAPE_TYPE, "AUTO_SHAPE", None)
PICTURE_SHAPE_TYPE = getattr(MSO_SHAPE_TYPE, "PICTURE", None)


def fail(message: str) -> None:
    raise SystemExit(message)


def check_internal_workflow_text(text: str, location: str) -> None:
    for pattern in INTERNAL_WORKFLOW_TEXT_PATTERNS:
        if pattern.search(text):
            fail(
                f"{location}: internal workflow/source-log text appears in "
                f"audience-facing content: {text[:100]}"
            )


def inches(value: int) -> float:
    return float(value) / EMU_PER_INCH


def visible_text(shape) -> str:
    if hasattr(shape, "text") and shape.text:
        return shape.text.strip()
    return ""


def words_in(text: str) -> list[str]:
    return re.findall(r"\b[\w'-]+\b", text)


def is_quiet_text(text: str) -> bool:
    stripped = text.strip()
    lower = stripped.lower()
    if lower.startswith(("source:", "sources:", "note:", "notes:", "figure:", "fig.")):
        return True
    if re.fullmatch(r"\d+(\s*/\s*\d+)?", stripped):
        return True
    return False


def looks_like_short_label(text: str) -> bool:
    stripped = " ".join(text.split())
    words = words_in(stripped)
    if len(words) > 4:
        return False
    return stripped.isupper() or len(stripped) <= 22


def text_shape_infos(slide) -> list[dict[str, object]]:
    infos: list[dict[str, object]] = []
    for shape in slide.shapes:
        text = visible_text(shape)
        if not text:
            continue
        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        infos.append(
            {
                "shape": shape,
                "text": text,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "right": left + width,
                "bottom": top + height,
                "words": words_in(text),
            }
        )
    return infos


def shape_area(shape) -> int:
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return max(width, 0) * max(height, 0)


def count_picture_shapes(shapes) -> int:
    count = 0
    for shape in shapes:
        if PICTURE_SHAPE_TYPE is not None and getattr(shape, "shape_type", None) == PICTURE_SHAPE_TYPE:
            count += 1
        elif hasattr(shape, "shapes"):
            count += count_picture_shapes(shape.shapes)
    return count


def shape_bounds(shape) -> tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return left, top, left + width, top + height


def info_center_inside(info: dict[str, object], shape, *, tolerance_in: float = 0.06) -> bool:
    tolerance = int(tolerance_in * EMU_PER_INCH)
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    right = left + int(getattr(shape, "width", 0) or 0)
    bottom = top + int(getattr(shape, "height", 0) or 0)
    cx = int(info["left"]) + int(info["width"]) / 2
    cy = int(info["top"]) + int(info["height"]) / 2
    return (
        left - tolerance <= cx <= right + tolerance
        and top - tolerance <= cy <= bottom + tolerance
    )


def is_container_candidate(shape, slide_area: int) -> bool:
    if visible_text(shape):
        return False
    if getattr(shape, "shape_type", None) in NON_EMPTY_VISUAL_TYPES:
        return False
    area = shape_area(shape)
    if area <= 0:
        return False
    area_ratio = area / slide_area
    if area_ratio < 0.018 or area_ratio > 0.45:
        return False
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return inches(width) >= 1.0 and inches(height) >= 0.45


def is_text_container_candidate(shape, slide_area: int) -> bool:
    if not visible_text(shape):
        return False
    if AUTO_SHAPE_TYPE is not None and getattr(shape, "shape_type", None) != AUTO_SHAPE_TYPE:
        return False
    area = shape_area(shape)
    if area <= 0:
        return False
    area_ratio = area / slide_area
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return 0.018 <= area_ratio <= 0.35 and inches(width) >= 1.0 and inches(height) >= 0.45


def run_font_size_pt(run, fallback: float) -> float:
    if run.font.size is not None:
        return float(run.font.size.pt)
    return fallback


def paragraph_font_size_pt(paragraph, fallback: float = 14.0) -> float:
    sizes = [
        run_font_size_pt(run, fallback)
        for run in paragraph.runs
        if run.text.strip()
    ]
    return max(sizes) if sizes else fallback


def text_frame_margins(shape) -> tuple[int, int, int, int]:
    if not hasattr(shape, "text_frame"):
        return (0, 0, 0, 0)
    tf = shape.text_frame
    return (
        int(tf.margin_left or 0),
        int(tf.margin_right or 0),
        int(tf.margin_top or 0),
        int(tf.margin_bottom or 0),
    )


def shape_max_font_size_pt(shape, fallback: float = 14.0) -> float:
    if not hasattr(shape, "text_frame"):
        return fallback
    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                sizes.append(run_font_size_pt(run, fallback))
    return max(sizes) if sizes else fallback


def shape_text_line_count(shape) -> int:
    text = visible_text(shape)
    if not text:
        return 0
    lines = [line for line in text.splitlines() if line.strip()]
    if lines:
        return len(lines)
    return 1


def estimate_text_required_height(shape, *, default_font_size: float) -> int:
    if not hasattr(shape, "text_frame"):
        return 0

    margin_left, margin_right, margin_top, margin_bottom = text_frame_margins(shape)
    usable_width = max(int(getattr(shape, "width", 0) or 0) - margin_left - margin_right, 1)
    usable_width_pt = usable_width / 12700
    required_pt = 0.0

    for paragraph in shape.text_frame.paragraphs:
        paragraph_text = "".join(run.text for run in paragraph.runs).strip()
        font_size = paragraph_font_size_pt(paragraph, fallback=default_font_size)
        if not paragraph_text:
            required_pt += font_size * 0.55
            continue

        avg_char_width = font_size * 0.52
        if re.search(r"[\u4e00-\u9fff]", paragraph_text):
            avg_char_width = font_size * 0.92
        chars_per_line = max(int(usable_width_pt / max(avg_char_width, 1)), 1)
        explicit_lines = [line for line in paragraph_text.splitlines() if line.strip()]
        line_count = 0
        for line in explicit_lines or [paragraph_text]:
            line_count += max(1, math.ceil(len(line) / chars_per_line))
        required_pt += line_count * font_size * 1.22

    return int(required_pt * 12700) + margin_top + margin_bottom


def combined_slide_text(slide) -> str:
    return "\n".join(
        visible_text(shape)
        for shape in slide.shapes
        if visible_text(shape)
    )


def is_closing_slide(slide) -> bool:
    text = combined_slide_text(slide).lower()
    has_thanks = bool(re.search(r"\b(thank\s*you|thanks)\b|谢谢|感谢", text))
    has_questions = bool(re.search(r"\bquestions?\b|问题|提问", text))
    return has_thanks and has_questions


def closing_has_only_thank_you_and_questions(slide) -> bool:
    text = "\n".join(
        visible_text(shape)
        for shape in slide.shapes
        if visible_text(shape) and not is_quiet_text(visible_text(shape))
    ).lower()
    stripped = re.sub(r"\b(thank\s*you|thanks)\b|谢谢|感谢", "", text)
    stripped = re.sub(r"\bquestions?\b|问题|提问", "", stripped)
    stripped = re.sub(r"\bdiscussion\b|讨论", "", stripped)
    stripped = re.sub(r"\band\b", "", stripped)
    stripped = re.sub(r"[^a-z0-9\u4e00-\u9fff]+", "", stripped)
    return stripped == ""


def cover_contains_duration(slide) -> bool:
    text = combined_slide_text(slide).lower()
    return bool(
        re.search(
            r"时长|持续时间|演讲时长|汇报时长|报告时长|"
            r"\bduration\b|\blength\b|\btime\s*limit\b|"
            r"\b\d+\s*(?:min|mins|minute|minutes|hour|hours)\b|"
            r"\d+\s*(?:分钟|小时)",
            text,
        )
    )


def cover_has_required_metadata(slide) -> bool:
    text = combined_slide_text(slide)
    lower = text.lower()
    presenter = bool(
        re.search(
            r"汇报人|报告人|主讲人|presenter|presented\s+by|speaker|prepared\s+by|name\s*:",
            lower,
        )
    )
    presentation_date = bool(
        re.search(
            r"\b20\d{2}(?:[-/.年]\d{1,2}(?:[-/.月]\d{1,2})?)?\b|"
            r"\b\d{1,2}[-/.]\d{1,2}(?:[-/.]\d{2,4})?\b|"
            r"\b(?:jan|feb|mar|apr|may|jun|jul|aug|sep|sept|oct|nov|dec)[a-z]*\b",
            lower,
        )
    )
    return presenter and presentation_date


def non_quiet_text_count(slide) -> int:
    return sum(
        1
        for shape in slide.shapes
        if visible_text(shape) and not is_quiet_text(visible_text(shape))
    )


def check_deck_shell(prs: Presentation, *, allow_missing_cover_closing: bool) -> None:
    if allow_missing_cover_closing:
        return
    if len(prs.slides) < 2:
        fail("deck must include a simple cover slide and a final Thank You + Questions slide")

    first = prs.slides[0]
    last = prs.slides[-1]

    if is_closing_slide(first):
        fail("slide 1 appears to be a closing slide; add a simple metadata cover first")
    if cover_contains_duration(first):
        fail("slide 1 cover must use presentation date, not presentation duration")
    if not cover_has_required_metadata(first):
        fail(
            "slide 1 must be a simple cover with title plus presentation date and presenter"
        )
    if non_quiet_text_count(first) > 8:
        fail(
            "slide 1 cover looks overloaded; keep only title, presentation date, presenter, "
            "and required basic metadata"
        )
    if not is_closing_slide(last):
        fail("final slide must contain Thank You plus Questions or Questions & Discussion")
    if not closing_has_only_thank_you_and_questions(last):
        fail("final slide must contain only Thank You plus Questions or Questions & Discussion")


def check_short_label_wrapping(slide_number: int, info: dict[str, object]) -> None:
    text = str(info["text"])
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) < 2:
        return
    joined = " ".join(lines)
    if len(words_in(joined)) > 4:
        return
    if inches(int(info["width"])) >= 2.3:
        return
    fail(
        f"slide {slide_number}: short label wraps across lines; widen or redesign it: "
        f"{joined[:60]}"
    )


def significant_vertical_intervals(slide, slide_area: int, slide_height: int) -> list[tuple[int, int]]:
    intervals: list[tuple[int, int]] = []
    for shape in slide.shapes:
        left, top, right, bottom = shape_bounds(shape)
        height = bottom - top
        area = shape_area(shape)
        if height <= 0 or area <= 0:
            continue
        area_ratio = area / slide_area
        if area_ratio > 0.85:
            continue
        if top < 0.35 * EMU_PER_INCH and height < 0.35 * EMU_PER_INCH:
            continue
        text = visible_text(shape)
        if text:
            if is_quiet_text(text):
                continue
            intervals.append((top, bottom))
            continue
        if area_ratio >= 0.012 and height >= 0.12 * EMU_PER_INCH:
            if bottom > slide_height * 0.94 and height < 0.35 * EMU_PER_INCH:
                continue
            intervals.append((top, bottom))
    return intervals


def merged_intervals(intervals: list[tuple[int, int]], *, join_gap_in: float = 0.08) -> list[tuple[int, int]]:
    if not intervals:
        return []
    join_gap = int(join_gap_in * EMU_PER_INCH)
    merged: list[tuple[int, int]] = []
    for start, end in sorted(intervals):
        if not merged or start > merged[-1][1] + join_gap:
            merged.append((start, end))
        else:
            merged[-1] = (merged[-1][0], max(merged[-1][1], end))
    return merged


def content_slide_range(prs: Presentation, *, skip_shell: bool) -> range:
    start = 1 if skip_shell and len(prs.slides) > 2 else 0
    end = len(prs.slides) - 1 if skip_shell and len(prs.slides) > 2 else len(prs.slides)
    return range(start, max(start, end))


def check_vertical_rhythm(
    prs: Presentation,
    *,
    max_vertical_gap: float,
    allow_large_vertical_gaps: bool,
    skip_shell: bool,
) -> None:
    if allow_large_vertical_gaps:
        return

    slide_height = int(prs.slide_height)
    slide_area = int(prs.slide_width) * int(prs.slide_height)
    max_gap = int(max_vertical_gap * EMU_PER_INCH)

    for slide_index in content_slide_range(prs, skip_shell=skip_shell):
        slide = prs.slides[slide_index]
        intervals = merged_intervals(
            significant_vertical_intervals(slide, slide_area, slide_height)
        )
        if len(intervals) < 3:
            continue
        for (_, prev_bottom), (next_top, _) in zip(intervals, intervals[1:]):
            gap = next_top - prev_bottom
            if gap <= max_gap:
                continue
            if prev_bottom < slide_height * 0.24:
                continue
            if next_top > slide_height * 0.92:
                continue
            fail(
                f"slide {slide_index + 1}: large empty vertical gap "
                f"({inches(gap):.2f} in) disrupts layout rhythm"
            )


def check_title_spacing(
    prs: Presentation,
    *,
    min_title_subtitle_gap: float,
    allow_tight_title_spacing: bool,
    skip_shell: bool,
) -> None:
    if allow_tight_title_spacing:
        return

    min_gap = int(min_title_subtitle_gap * EMU_PER_INCH)
    slide_height = int(prs.slide_height)

    for slide_index in content_slide_range(prs, skip_shell=skip_shell):
        slide = prs.slides[slide_index]
        candidates = []
        for shape in slide.shapes:
            text = visible_text(shape)
            if not text or is_quiet_text(text):
                continue
            font_size = shape_max_font_size_pt(shape)
            left, top, right, bottom = shape_bounds(shape)
            if top <= slide_height * 0.34 and font_size >= 26:
                candidates.append((top, -font_size, shape))
        if not candidates:
            continue

        _, _, title_shape = sorted(candidates)[0]
        title_left, title_top, title_right, title_bottom = shape_bounds(title_shape)
        title_font = shape_max_font_size_pt(title_shape)
        estimated_title_height = estimate_text_required_height(
            title_shape,
            default_font_size=title_font,
        )
        title_content_bottom = title_top + min(
            max(estimated_title_height, 0),
            title_bottom - title_top,
        )
        title_lines = shape_text_line_count(title_shape)
        if title_lines >= 2:
            actual_height_pt = (title_bottom - title_top) / 12700
            required_height_pt = title_lines * title_font * 1.08
            if actual_height_pt < required_height_pt:
                fail(
                    f"slide {slide_index + 1}: multi-line title spacing is too tight; "
                    "increase title box height or line spacing"
                )

        subtitle_candidates = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            text = visible_text(shape)
            if not text or is_quiet_text(text):
                continue
            font_size = shape_max_font_size_pt(shape)
            left, top, right, bottom = shape_bounds(shape)
            horizontal_overlap = min(title_right, right) - max(title_left, left)
            aligned_with_title = left <= title_left + 0.45 * EMU_PER_INCH or horizontal_overlap > 0
            if (
                aligned_with_title
                and title_top < top < slide_height * 0.42
                and font_size < title_font
            ):
                subtitle_candidates.append((top, id(shape), shape))
        if not subtitle_candidates:
            continue
        subtitle_top, _, _ = sorted(subtitle_candidates)[0]
        gap = subtitle_top - title_content_bottom
        if gap < min_gap:
            fail(
                f"slide {slide_index + 1}: title and subtitle are too close "
                f"({inches(max(gap, 0)):.2f} in gap)"
            )


def check_text_box_capacity(
    prs: Presentation,
    *,
    default_font_size: float,
    overflow_tolerance: float,
    allow_estimated_text_overflow: bool,
) -> None:
    if allow_estimated_text_overflow:
        return

    slide_area = int(prs.slide_width) * int(prs.slide_height)

    for slide_number, slide in enumerate(prs.slides, start=1):
        cards = [
            shape
            for shape in slide.shapes
            if is_container_candidate(shape, slide_area)
        ]
        for shape in slide.shapes:
            text = visible_text(shape)
            if not text or is_quiet_text(text):
                continue
            if looks_like_short_label(text):
                continue
            info = {
                "left": int(getattr(shape, "left", 0) or 0),
                "top": int(getattr(shape, "top", 0) or 0),
                "width": int(getattr(shape, "width", 0) or 0),
                "height": int(getattr(shape, "height", 0) or 0),
            }
            info["right"] = int(info["left"]) + int(info["width"])
            info["bottom"] = int(info["top"]) + int(info["height"])
            inside_card = any(info_center_inside(info, card, tolerance_in=0.02) for card in cards)
            if not inside_card and not is_text_container_candidate(shape, slide_area):
                continue
            required_height = estimate_text_required_height(
                shape,
                default_font_size=default_font_size,
            )
            actual_height = int(getattr(shape, "height", 0) or 0)
            if required_height <= 0 or actual_height <= 0:
                continue
            if required_height > actual_height * overflow_tolerance:
                fail(
                    f"slide {slide_number}: text may overflow its text box "
                    f"(estimated {inches(required_height):.2f} in needed > "
                    f"{inches(actual_height):.2f} in available): {text[:80]}"
                )


def check_card_text_safety(
    prs: Presentation,
    *,
    min_card_text_padding: float,
    allow_card_edge_overflow: bool,
) -> None:
    if allow_card_edge_overflow:
        return

    padding = int(min_card_text_padding * EMU_PER_INCH)
    slide_area = int(prs.slide_width) * int(prs.slide_height)

    for slide_number, slide in enumerate(prs.slides, start=1):
        cards = [
            shape
            for shape in slide.shapes
            if is_container_candidate(shape, slide_area)
        ]
        if not cards:
            continue
        for info in text_shape_infos(slide):
            text = str(info["text"])
            if is_quiet_text(text):
                continue
            containing_cards = [
                shape for shape in cards if info_center_inside(info, shape, tolerance_in=0.02)
            ]
            if not containing_cards:
                continue
            card = min(containing_cards, key=shape_area)
            left, top, right, bottom = shape_bounds(card)
            text_left = int(info["left"])
            text_top = int(info["top"])
            text_right = int(info["right"])
            text_bottom = int(info["bottom"])
            if (
                text_left < left + padding
                or text_top < top + padding
                or text_right > right - padding
                or text_bottom > bottom - padding
            ):
                fail(
                    f"slide {slide_number}: text box is outside the card safe area; "
                    f"resize the card/text box or add padding: {text[:80]}"
                )


def check_large_sparse_containers(
    prs: Presentation,
    *,
    max_large_empty_shape_area_ratio: float,
    allow_large_empty_shapes: bool,
) -> None:
    slide_area = int(prs.slide_width) * int(prs.slide_height)

    for slide_number, slide in enumerate(prs.slides, start=1):
        infos = text_shape_infos(slide)
        for info in infos:
            check_short_label_wrapping(slide_number, info)

        for shape in slide.shapes:
            area = shape_area(shape)
            if area <= 0:
                continue
            area_ratio = area / slide_area
            if area_ratio < max_large_empty_shape_area_ratio:
                continue
            if area_ratio > 0.85:
                continue
            if inches(int(getattr(shape, "height", 0) or 0)) < 1.4:
                continue
            if getattr(shape, "shape_type", None) in NON_EMPTY_VISUAL_TYPES:
                continue

            text = visible_text(shape)
            if text:
                if (
                    not allow_large_empty_shapes
                    and not is_quiet_text(text)
                    and inches(int(getattr(shape, "height", 0) or 0)) >= 2.4
                    and len(words_in(text)) < 28
                ):
                    fail(
                        f"slide {slide_number}: large text container has too little content "
                        f"for its size: {text[:80]}"
                    )
                continue

            inside = [info for info in infos if info_center_inside(info, shape)]
            if not inside:
                if not allow_large_empty_shapes:
                    fail(
                        f"slide {slide_number}: large empty card/panel detected "
                        f"({area_ratio:.0%} of slide). Shrink it or add structured content."
                    )
                continue

            if allow_large_empty_shapes:
                continue

            top = int(getattr(shape, "top", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            relative_text_bottom = max((int(info["bottom"]) - top) / height for info in inside)
            inside_text_area = sum(int(info["width"]) * int(info["height"]) for info in inside)
            inside_text_area_ratio = inside_text_area / area
            inside_word_count = sum(len(info["words"]) for info in inside)

            if relative_text_bottom < 0.45 and (
                inside_text_area_ratio < 0.16 or inside_word_count < 35
            ):
                fail(
                    f"slide {slide_number}: large card/panel is visually sparse; "
                    "content sits near the top while most of the frame is empty."
                )


def check_slide_balance(prs: Presentation, *, allow_top_heavy: bool, skip_shell: bool) -> None:
    if allow_top_heavy:
        return

    slide_height = int(prs.slide_height)
    slide_area = int(prs.slide_width) * int(prs.slide_height)
    lower_threshold = slide_height * 0.68

    for slide_index in content_slide_range(prs, skip_shell=skip_shell):
        slide = prs.slides[slide_index]
        slide_number = slide_index + 1
        infos = [
            info
            for info in text_shape_infos(slide)
            if not is_quiet_text(str(info["text"]))
        ]
        if len(infos) < 3:
            continue

        max_text_bottom = max(int(info["bottom"]) for info in infos)
        if max_text_bottom >= lower_threshold:
            continue

        lower_visual = False
        for shape in slide.shapes:
            area = shape_area(shape)
            if area / slide_area < 0.012:
                continue
            bottom = int(getattr(shape, "top", 0) or 0) + int(getattr(shape, "height", 0) or 0)
            if bottom >= lower_threshold and area / slide_area < 0.85:
                lower_visual = True
                break

        if not lower_visual:
            fail(
                f"slide {slide_number}: meaningful content is top-heavy; "
                "use the lower-middle area or redesign the layout balance."
            )


def collect_slide_text_and_fonts(
    prs: Presentation,
    *,
    strict: bool = False,
    max_textbox_words: int = 35,
    min_body_font_size: float = 11.0,
) -> tuple[list[str], set[str]]:
    slide_text: list[str] = []
    fonts: set[str] = set()

    for slide_number, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            shape_text = visible_text(shape)
            if shape_text:
                texts.append(shape_text)
                if strict:
                    upper = shape_text.upper()
                    for marker in PLACEHOLDER_PATTERNS:
                        if marker in upper:
                            fail(f"slide {slide_number}: placeholder text found: {marker}")
                    check_internal_workflow_text(shape_text, f"slide {slide_number}")
                    words = re.findall(r"\b[\w'-]+\b", shape_text)
                    is_quiet_source = shape_text.lower().startswith(("source:", "sources:"))
                    if len(words) > max_textbox_words and not is_quiet_source:
                        fail(
                            f"slide {slide_number}: text box is too dense "
                            f"({len(words)} words > {max_textbox_words}): {shape_text[:80]}"
                        )
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        if (
                            strict
                            and run.text.strip()
                            and run.font.size is not None
                            and run.font.size.pt < min_body_font_size
                            and not is_quiet_text(shape_text)
                            and not looks_like_short_label(shape_text)
                        ):
                            fail(
                                f"slide {slide_number}: body text is too small "
                                f"({run.font.size.pt:.1f} pt < {min_body_font_size:.1f} pt): "
                                f"{shape_text[:80]}"
                            )

            if shape.left < 0 or shape.top < 0:
                fail(f"slide {slide_number}: shape outside top/left boundary")
            if shape.left + shape.width > prs.slide_width:
                fail(f"slide {slide_number}: shape outside right boundary")
            if shape.top + shape.height > prs.slide_height:
                fail(f"slide {slide_number}: shape outside bottom boundary")

        if not texts:
            fail(f"slide {slide_number}: no visible text found")
        slide_text.append("\n".join(texts))

    return slide_text, fonts


def scan_forbidden_colors(pptx: Path) -> set[str]:
    found: set[str] = set()
    with zipfile.ZipFile(pptx) as zf:
        for name in zf.namelist():
            if not name.startswith("ppt/slides/") or not name.endswith(".xml"):
                continue
            text = zf.read(name).decode("utf-8", errors="ignore").upper()
            for color in FORBIDDEN_COLOR_HEXES:
                if color in text:
                    found.add(color)
    return found


def validate_script_markers(script: Path, slide_count: int) -> None:
    try:
        from docx import Document
    except Exception as exc:  # pragma: no cover
        fail(f"python-docx is required to inspect script markers: {exc}")

    text = "\n".join(p.text for p in Document(script).paragraphs)
    check_internal_workflow_text(text, "script")
    for index in range(1, slide_count + 1):
        if not re.search(rf"\bSlide\s+{index}\b", text, flags=re.IGNORECASE):
            fail(f"script missing Slide {index} marker")


def render_previews(pptx: Path, out_dir: Path) -> list[Path]:
    qlmanage = shutil.which("qlmanage")
    if qlmanage is None:
        fail("qlmanage not found; cannot render previews with --render")

    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(pptx)
    split_paths: list[Path] = []

    for keep in range(len(prs.slides)):
        one = Presentation(pptx)
        sld_id_lst = one.slides._sldIdLst
        slides = list(sld_id_lst)
        for idx in reversed(range(len(slides))):
            if idx != keep:
                rel_id = slides[idx].rId
                one.part.drop_rel(rel_id)
                sld_id_lst.remove(slides[idx])
        split_path = out_dir / f"slide_{keep + 1}.pptx"
        one.save(split_path)
        split_paths.append(split_path)

    png_paths: list[Path] = []
    for split_path in split_paths:
        subprocess.run(
            [qlmanage, "-t", "-s", "1400", "-o", str(out_dir), str(split_path)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        png = out_dir / f"{split_path.name}.png"
        if not png.exists() or png.stat().st_size == 0:
            fail(f"preview render failed for {split_path.name}")
        png_paths.append(png)

    return png_paths


def main() -> None:
    parser = argparse.ArgumentParser(description="Audit a PPTPRO-style deck.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--script", type=Path)
    parser.add_argument("--render", action="store_true")
    parser.add_argument("--preview-dir", type=Path)
    parser.add_argument("--allow-serif", action="store_true")
    parser.add_argument("--allow-beige", action="store_true")
    parser.add_argument(
        "--strict",
        action="store_true",
        help="Require rendered previews and stricter text-density checks.",
    )
    parser.add_argument(
        "--max-textbox-words",
        type=int,
        default=35,
        help="Maximum words allowed in a single non-source text box under --strict.",
    )
    parser.add_argument(
        "--min-pictures",
        type=int,
        default=0,
        help="Minimum picture count required, useful after inventorying source document figures.",
    )
    parser.add_argument(
        "--min-body-font-size",
        type=float,
        default=11.0,
        help="Minimum explicit body font size under --strict; footers and labels are ignored.",
    )
    parser.add_argument(
        "--default-font-size",
        type=float,
        default=14.0,
        help="Fallback font size used for strict text overflow estimates.",
    )
    parser.add_argument(
        "--text-overflow-tolerance",
        type=float,
        default=1.80,
        help="Allowed estimated text-height ratio before strict overflow fails.",
    )
    parser.add_argument(
        "--min-card-text-padding",
        type=float,
        default=0.10,
        help="Minimum text inset in inches when a text box sits inside a detected card/panel.",
    )
    parser.add_argument(
        "--max-large-empty-shape-area-ratio",
        type=float,
        default=0.12,
        help="Largest allowed sparse card/panel area ratio under --strict.",
    )
    parser.add_argument(
        "--allow-large-empty-shapes",
        action="store_true",
        help="Allow large sparse cards/panels under --strict.",
    )
    parser.add_argument(
        "--allow-top-heavy",
        action="store_true",
        help="Allow slides whose meaningful content stays in the upper part of the canvas.",
    )
    parser.add_argument(
        "--max-vertical-gap",
        type=float,
        default=1.45,
        help="Maximum empty vertical gap in inches between significant content bands under --strict.",
    )
    parser.add_argument(
        "--min-title-subtitle-gap",
        type=float,
        default=0.12,
        help="Minimum title-to-subtitle gap in inches under --strict.",
    )
    parser.add_argument(
        "--allow-large-vertical-gaps",
        action="store_true",
        help="Allow large empty vertical gaps between significant slide content bands.",
    )
    parser.add_argument(
        "--allow-tight-title-spacing",
        action="store_true",
        help="Allow cramped title line spacing or tight title/subtitle spacing.",
    )
    parser.add_argument(
        "--allow-missing-cover-closing",
        action="store_true",
        help="Allow decks without the default simple cover and final Thank You + Questions slide.",
    )
    parser.add_argument(
        "--allow-card-edge-overflow",
        action="store_true",
        help="Allow text boxes to touch or exceed detected card/panel safe areas.",
    )
    parser.add_argument(
        "--allow-estimated-text-overflow",
        action="store_true",
        help="Disable estimated text box capacity checks under --strict.",
    )
    args = parser.parse_args()

    pptx = args.pptx.expanduser().resolve()
    if not pptx.exists():
        fail(f"PPTX not found: {pptx}")

    prs = Presentation(pptx)
    if len(prs.slides) == 0:
        fail("deck has no slides")

    if args.strict and not args.render:
        fail("--strict requires --render so visual previews are generated")

    picture_count = sum(count_picture_shapes(slide.shapes) for slide in prs.slides)
    if picture_count < args.min_pictures:
        fail(
            f"deck contains too few picture assets "
            f"({picture_count} found < {args.min_pictures} required)"
        )

    slide_text, fonts = collect_slide_text_and_fonts(
        prs,
        strict=args.strict,
        max_textbox_words=args.max_textbox_words,
        min_body_font_size=args.min_body_font_size,
    )

    if args.strict:
        check_deck_shell(
            prs,
            allow_missing_cover_closing=args.allow_missing_cover_closing,
        )
        skip_shell = not args.allow_missing_cover_closing
        check_text_box_capacity(
            prs,
            default_font_size=args.default_font_size,
            overflow_tolerance=args.text_overflow_tolerance,
            allow_estimated_text_overflow=args.allow_estimated_text_overflow,
        )
        check_card_text_safety(
            prs,
            min_card_text_padding=args.min_card_text_padding,
            allow_card_edge_overflow=args.allow_card_edge_overflow,
        )
        check_large_sparse_containers(
            prs,
            max_large_empty_shape_area_ratio=args.max_large_empty_shape_area_ratio,
            allow_large_empty_shapes=args.allow_large_empty_shapes,
        )
        check_vertical_rhythm(
            prs,
            max_vertical_gap=args.max_vertical_gap,
            allow_large_vertical_gaps=args.allow_large_vertical_gaps,
            skip_shell=skip_shell,
        )
        check_title_spacing(
            prs,
            min_title_subtitle_gap=args.min_title_subtitle_gap,
            allow_tight_title_spacing=args.allow_tight_title_spacing,
            skip_shell=skip_shell,
        )
        check_slide_balance(
            prs,
            allow_top_heavy=args.allow_top_heavy,
            skip_shell=skip_shell,
        )

    if not args.allow_serif:
        bad_fonts = fonts & FORBIDDEN_FONTS
        if bad_fonts:
            fail(f"forbidden serif fonts found: {sorted(bad_fonts)}")

    if not args.allow_beige:
        bad_colors = scan_forbidden_colors(pptx)
        if bad_colors:
            fail(f"forbidden beige-like colors found: {sorted(bad_colors)}")

    if args.script:
        validate_script_markers(args.script.expanduser().resolve(), len(prs.slides))

    preview_paths: list[Path] = []
    if args.render:
        preview_dir = args.preview_dir
        if preview_dir is None:
            preview_dir = pptx.parent / "pptpro_previews"
        preview_paths = render_previews(pptx, preview_dir.expanduser().resolve())

    print(
        f"PPTPRO audit OK | slides={len(slide_text)} | "
        f"pictures={picture_count} | fonts={sorted(fonts) or ['theme/default']}"
    )
    if preview_paths:
        for path in preview_paths:
            print(path)


if __name__ == "__main__":
    main()
